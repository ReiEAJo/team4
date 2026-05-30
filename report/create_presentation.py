"""
이 모듈은 1인가구 상권 분석 보고서 내용(EDA 및 대시보드 결과)을 바탕으로 
파이썬 python-pptx 라이브러리를 활용해 프레젠테이션 파일(presentation.pptx)을 자동 생성하는 역할을 수행합니다.
주요 기능:
- 슬라이드 레이아웃 및 테마 색상 설정
- 표지, 목차, 간지, 본문 슬라이드 구성
- python-pptx 내장 차트(Bar, Pie, Column)를 활용한 시각화
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.chart import XL_LEGEND_POSITION

def create_presentation():
    # 1. Create a presentation object
    prs = Presentation()
    
    # 16:9 aspect ratio
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # Color Palette
    primary_color = RGBColor(30, 39, 97) # 1E2761 Navy
    accent_gold = RGBColor(249, 231, 149) # F9E795 Gold
    accent_coral = RGBColor(249, 97, 103) # F96167 Coral
    text_dark = RGBColor(54, 54, 54) # 363636 Dark Gray
    text_light = RGBColor(255, 255, 255) # White
    bg_light = RGBColor(245, 245, 245) # F5F5F5 Off-white
    
    # Layouts
    blank_slide_layout = prs.slide_layouts[6]

    def add_shape(slide, shape_type, left, top, width, height, fill_color, line_color=None):
        shape = slide.shapes.add_shape(shape_type, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        if line_color:
            shape.line.color.rgb = line_color
        else:
            shape.line.fill.background()
        return shape

    def add_text(slide, text, left, top, width, height, font_size, font_color, bold=False, align=PP_ALIGN.LEFT):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.font.name = 'Malgun Gothic'
        p.alignment = align
        return txBox

    # --- Slide 1: Cover ---
    slide1 = prs.slides.add_slide(blank_slide_layout)
    add_shape(slide1, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(5.625), primary_color)
    add_shape(slide1, MSO_SHAPE.RECTANGLE, Inches(1), Inches(4), Inches(8), Inches(0.05), accent_gold)
    
    add_text(slide1, "서울시 1인가구 상권 분석 및 창업 입지 제안", Inches(1), Inches(1.5), Inches(8), Inches(1), 32, text_light, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide1, "유동인구 대비 실매출 분석을 중심으로", Inches(1), Inches(2.5), Inches(8), Inches(0.5), 20, accent_gold, align=PP_ALIGN.CENTER)
    add_text(slide1, "팀명: [TBD]\n팀원: [TBD]\n날짜: 2026. 05. 30", Inches(1), Inches(4.2), Inches(8), Inches(1), 14, text_light, align=PP_ALIGN.CENTER)

    # --- Slide 2: Table of Contents ---
    slide2 = prs.slides.add_slide(blank_slide_layout)
    add_shape(slide2, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(5.625), bg_light)
    add_text(slide2, "목차 (Contents)", Inches(0.5), Inches(0.5), Inches(9), Inches(0.8), 28, primary_color, bold=True)
    
    contents = [
        "1. 문제 정의",
        "2. 분석 배경 및 목적",
        "3. 데이터 소개 (데이터 개요 및 전처리 요약)",
        "4. 분석 결과 및 인사이트",
        "5. 결론 및 비즈니스 액션플랜"
    ]
    for i, c in enumerate(contents):
        add_text(slide2, c, Inches(1), Inches(1.5 + i * 0.7), Inches(8), Inches(0.5), 18, text_dark, bold=True)

    def create_divider_slide(title, subtitle):
        slide = prs.slides.add_slide(blank_slide_layout)
        add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(5.625), primary_color)
        add_text(slide, title, Inches(1), Inches(2.2), Inches(8), Inches(0.8), 36, text_light, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, subtitle, Inches(1), Inches(3.2), Inches(8), Inches(0.5), 20, accent_coral, bold=True, align=PP_ALIGN.CENTER)
        return slide

    def create_content_slide(title):
        slide = prs.slides.add_slide(blank_slide_layout)
        add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(5.625), text_light)
        # Header bar
        add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.2), Inches(5.625), primary_color)
        add_text(slide, title, Inches(0.5), Inches(0.4), Inches(9), Inches(0.6), 24, primary_color, bold=True)
        return slide

    # --- Slide 3: Divider 1 ---
    create_divider_slide("Part 1. 문제 정의", "Problem Definition")

    # --- Slide 4: Background & Objective ---
    slide4 = create_content_slide("분석 배경 및 목적")
    bg_obj = (
        "분석 배경\n"
        "• 1인가구 비중 급증으로 상권 구조 변화\n"
        "• 유동인구가 많은 상권이 반드시 생존율이 높지 않다는 '유동인구의 함정' 가설 대두\n\n"
        "분석 목적\n"
        "• 1인가구 밀집 지역 데이터를 기반으로 가장 생존율이 높은(낮은 폐업률) 창업 입지 도출\n"
        "• 해당 상권 내 가장 유리한 비즈니스 아이템(업종) 제안"
    )
    add_text(slide4, bg_obj, Inches(0.5), Inches(1.5), Inches(9), Inches(3.5), 18, text_dark)

    # --- Slide 5: Divider 2 ---
    create_divider_slide("Part 2. 데이터 소개", "Data Introduction")

    # --- Slide 6: Data Overview ---
    slide6 = create_content_slide("데이터 개요")
    data_overview = (
        "• 데이터 출처: 서울시 우리마을가게 상권분석서비스 공공데이터\n"
        "• 수집 기간: 2022년 ~ 2023년\n"
        "• 분석 대상: 서울시 1인가구 밀집 주요 행정동 236개 상권\n"
        "• 데이터 크기: 236개 레코드, 5개 수치형 컬럼 및 파생변수\n\n"
        "주요 지표 분포 요약\n"
        "• 1인가구 분포는 상위 지역에 극도로 밀집 (최대 18,557명)\n"
        "• 폐업률은 평균 2.7% 수준이나, 최대 5.7%까지 치솟는 고위험 지역 존재"
    )
    add_text(slide6, data_overview, Inches(0.5), Inches(1.5), Inches(9), Inches(3.5), 16, text_dark)

    # --- Slide 7: Preprocessing ---
    slide7 = create_content_slide("데이터 전처리 요약")
    add_text(slide7, "신뢰도 확보를 위한 전처리 과정", Inches(0.5), Inches(1.5), Inches(9), Inches(0.5), 18, primary_color, bold=True)
    preprocessing = (
        "1. 결측치 처리\n"
        "   - 핵심 지표 누락이 있는 결측치는 분석 대상에서 제외하여 데이터 무결성 확보\n\n"
        "2. 이상치(Outlier) 보정\n"
        "   - 상위 5% 극단적 매출액 및 유동인구에 IQR(Interquartile Range) 방식 적용\n"
        "   - 1.5 × IQR 범위를 벗어나는 데이터를 보정하여 통계적 왜곡 방지\n\n"
        "3. 파생변수 생성\n"
        "   - 주말 매출 비율, 2030세대 비율, 여성 비율 등 분석용 파생변수 도출"
    )
    add_text(slide7, preprocessing, Inches(0.5), Inches(2.2), Inches(9), Inches(3), 16, text_dark)

    # --- Slide 8: Divider 3 ---
    create_divider_slide("Part 3. 분석 결과 및 인사이트", "Analysis Results & Insights")

    # --- Slide 9: 1인가구 타겟 (차트) ---
    slide9 = create_content_slide("배후 수요 파악: 자치구별 1인가구 밀집도")
    add_text(slide9, "관악구, 강남구, 강서구 등 상위 3개 자치구의 1인가구 수가 서울 평균 대비 2배 이상 높음.", Inches(0.5), Inches(1.2), Inches(9), Inches(0.5), 14, text_dark)
    
    chart_data = CategoryChartData()
    chart_data.categories = ['관악구', '강서구', '강남구', '동작구', '영등포구']
    chart_data.add_series('1인가구 총계', (186002, 128965, 97631, 91790, 98382))
    
    x, y, cx, cy = Inches(0.5), Inches(1.8), Inches(9), Inches(3.5)
    chart = slide9.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    ).chart
    chart.has_title = False
    chart.has_legend = False

    # --- Slide 10: 소비 패턴 ---
    slide10 = create_content_slide("1인가구 핵심 소비 카테고리 (Top 5)")
    add_text(slide10, "지역 내 소비액의 70% 이상이 '음식(외식 및 배달)'과 '식료품(편의점)'에 집중되어 있음.", Inches(0.5), Inches(1.2), Inches(9), Inches(0.5), 14, text_dark)
    
    chart_data_pie = CategoryChartData()
    chart_data_pie.categories = ['외식업(음식)', '소매업(식료품 등)', '서비스업(여가)', '기타']
    chart_data_pie.add_series('소비 비중', (45, 30, 15, 10))
    
    chart2 = slide10.shapes.add_chart(
        XL_CHART_TYPE.PIE, Inches(2), Inches(1.8), Inches(6), Inches(3.5), chart_data_pie
    ).chart
    chart2.has_legend = True
    chart2.legend.position = XL_LEGEND_POSITION.RIGHT
    chart2.plots[0].has_data_labels = True

    # --- Slide 11: 유동인구의 함정 ---
    slide11 = create_content_slide("유동인구의 함정: 사람이 많다고 돈을 버는가?")
    
    add_shape(slide11, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(4.2), Inches(3.5), bg_light)
    add_text(slide11, "상관관계 분석 결과", Inches(0.7), Inches(1.7), Inches(3.8), Inches(0.5), 18, primary_color, bold=True)
    add_text(slide11, "• 매출액과 폐업률: 강한 음의 상관관계 (r < -0.5)\n• 유동인구와 폐업률: 상관관계 미미", Inches(0.7), Inches(2.2), Inches(3.8), Inches(2), 16, text_dark)
    
    add_shape(slide11, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.3), Inches(1.5), Inches(4.2), Inches(3.5), bg_light, line_color=accent_coral)
    add_text(slide11, "핵심 인사이트", Inches(5.5), Inches(1.7), Inches(3.8), Inches(0.5), 18, accent_coral, bold=True)
    add_text(slide11, "유동인구가 많은 비싼 대로변 상권보다,\n\n실결제 매출액이 높으면서 임대료가 합리적인 '이면도로(골목)' 점포가 생존에 훨씬 유리합니다.", Inches(5.5), Inches(2.2), Inches(3.8), Inches(2.5), 16, text_dark)

    # --- Slide 12: 최적 입지 Top 5 ---
    slide12 = create_content_slide("추천 입지 Top 5 (1인가구 및 실매출 기반)")
    add_text(slide12, "매출 볼륨이 크고 폐업률이 낮아 안정성이 높은 1인가구 밀집 지역", Inches(0.5), Inches(1.2), Inches(9), Inches(0.5), 14, text_dark)
    
    # Table creation
    rows, cols = 6, 5
    left, top, width, height = Inches(0.5), Inches(1.8), Inches(9), Inches(3)
    table = slide12.shapes.add_table(rows, cols, left, top, width, height).table
    
    columns = ["순위", "행정동명", "1인가구 총계", "당월 매출액(만원)", "폐업률(%)"]
    data = [
        ["1", "서교동", "10,405", "654,812", "2.13%"],
        ["2", "역삼1동", "17,721", "928,476", "2.32%"],
        ["3", "문정2동", "8,963", "281,619", "2.33%"],
        ["4", "가산동", "14,731", "1,425,070", "2.86%"],
        ["5", "영등포동", "17,144", "482,078", "2.65%"]
    ]
    
    for c, col_name in enumerate(columns):
        table.cell(0, c).text = col_name
        table.cell(0, c).text_frame.paragraphs[0].font.bold = True
        table.cell(0, c).fill.solid()
        table.cell(0, c).fill.fore_color.rgb = primary_color
        table.cell(0, c).text_frame.paragraphs[0].font.color.rgb = text_light
        
    for r, row in enumerate(data):
        for c, val in enumerate(row):
            table.cell(r+1, c).text = val
            table.cell(r+1, c).text_frame.paragraphs[0].font.color.rgb = text_dark

    # --- Slide 13: 유망 창업 업종 ---
    slide13 = create_content_slide("추천 상권 내 최적의 창업 업종(아이템) 심층 분석")
    add_text(slide13, "해당 상권에서는 대형 시설보다 '1인 맞춤 생활밀착형 서비스업'의 폐업 리스크가 가장 낮게 나타납니다.", Inches(0.5), Inches(1.2), Inches(9), Inches(0.5), 14, text_dark)
    
    chart_data_bar = CategoryChartData()
    chart_data_bar.categories = ['네일숍', '반찬가게', '세탁소', '미용실', '슈퍼마켓', '피부관리실', '사진관']
    chart_data_bar.add_series('폐업률(%)', (1.65, 1.95, 2.50, 2.60, 2.70, 2.85, 2.95))
    
    chart3 = slide13.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.5), Inches(1.8), Inches(9), Inches(3.5), chart_data_bar
    ).chart
    chart3.has_legend = False
    chart3.plots[0].has_data_labels = True

    # --- Slide 14: Divider 4 ---
    create_divider_slide("Part 4. 결론 및 비즈니스 액션플랜", "Conclusion & Business Action Plan")

    # --- Slide 15: Key Findings ---
    slide15 = create_content_slide("인사이트 요약 (Key Findings)")
    add_shape(slide15, MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.5), Inches(9), Inches(1), bg_light)
    add_text(slide15, "1. 유동인구의 함정", Inches(0.7), Inches(1.6), Inches(8), Inches(0.3), 16, primary_color, bold=True)
    add_text(slide15, "유동인구 규모보다 실결제 매출액 방어가 상권 생존에 더 중요합니다.", Inches(0.7), Inches(2.0), Inches(8), Inches(0.3), 14, text_dark)
    
    add_shape(slide15, MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2.7), Inches(9), Inches(1), bg_light)
    add_text(slide15, "2. 추천 상권 (Top 5)", Inches(0.7), Inches(2.8), Inches(8), Inches(0.3), 16, primary_color, bold=True)
    add_text(slide15, "서교동, 역삼1동, 문정2동, 가산동, 영등포동이 안정적인 창업 지역으로 분석되었습니다.", Inches(0.7), Inches(3.2), Inches(8), Inches(0.3), 14, text_dark)

    add_shape(slide15, MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(3.9), Inches(9), Inches(1), bg_light)
    add_text(slide15, "3. 유망 창업 아이템", Inches(0.7), Inches(4.0), Inches(8), Inches(0.3), 16, primary_color, bold=True)
    add_text(slide15, "1인 맞춤 생활밀착형(반찬가게, 코인세탁소 등) 및 외식업이 폐업 리스크가 가장 낮습니다.", Inches(0.7), Inches(4.4), Inches(8), Inches(0.3), 14, text_dark)

    # --- Slide 16: Action Plan ---
    slide16 = create_content_slide("비즈니스 액션플랜")
    
    add_text(slide16, "단기 액션: 1인가구 특화 저위험 결합 모델 창업", Inches(0.5), Inches(1.3), Inches(9), Inches(0.4), 18, primary_color, bold=True)
    ap1 = (
        "• 전략: '무인 반찬 자판기 + 코인세탁소'가 결합된 멀티숍 창업\n"
        "• 입지: 서교동, 역삼1동 등 핵심 상권의 이면도로 (임대료 약 30% 저렴)\n"
        "• 실행 주체: 예비 소상공인 창업자\n"
        "• KPI: 초기 창업 비용 15% 절감, 창업 후 1년 생존율 90% 확보"
    )
    add_text(slide16, ap1, Inches(0.5), Inches(1.8), Inches(9), Inches(1.5), 14, text_dark)

    add_shape(slide16, MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(3.2), Inches(9), Inches(0.02), primary_color)
    
    add_text(slide16, "중장기 액션: 프랜차이즈 가맹본부의 1인 특화 상권 전략 출점", Inches(0.5), Inches(3.5), Inches(9), Inches(0.4), 18, accent_coral, bold=True)
    ap2 = (
        "• 전략: 마라탕, 소포장 보쌈 등 1인 타겟 외식 브랜드의 핵심 직영점 거점으로 가산동/서교동 선정\n"
        "• 액션: 오피스타운 런치 세트와 야간 배달 세트 이원화 운영\n"
        "• 실행 주체: F&B 프랜차이즈 점포개발팀\n"
        "• KPI: 타 상권 가맹점 대비 평균 월매출 25% 상회"
    )
    add_text(slide16, ap2, Inches(0.5), Inches(4.0), Inches(9), Inches(1.5), 14, text_dark)

    # --- Slide 17: Outro ---
    slide17 = prs.slides.add_slide(blank_slide_layout)
    add_shape(slide17, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(5.625), primary_color)
    add_text(slide17, "Thank You", Inches(1), Inches(2.2), Inches(8), Inches(0.8), 48, text_light, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide17, "Q & A", Inches(1), Inches(3.2), Inches(8), Inches(0.5), 24, accent_gold, bold=True, align=PP_ALIGN.CENTER)

    # Save presentation
    output_path = os.path.join(os.path.dirname(__file__), 'presentation.pptx')
    prs.save(output_path)
    print(f"Presentation generated successfully at: {output_path}")

if __name__ == "__main__":
    create_presentation()
